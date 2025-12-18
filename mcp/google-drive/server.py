from mcp.server.fastmcp import FastMCP
import os
import io
from google.oauth2.credentials import Credentials
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload, MediaIoBaseUpload
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
import openpyxl
import json
mcp = FastMCP("google-drive", host="0.0.0.0", port=8080)

import json

# ... imports ...

@mcp.tool()
def list_files(token: str, folder_id: str = 'root', order_by: str = 'folder,name') -> str:
    """
    List ALL files from a Google Drive folder. Returns a JSON string.
    
    Args:
        token: The OAuth2 access token for the user.
        folder_id: The ID of the folder to list files from (default: 'root').
        order_by: Sort order (e.g., 'folder,name', 'modifiedTime desc').
    """
    try:
        creds = Credentials(token=token)
        service = build('drive', 'v3', credentials=creds)

        all_files = []
        page_token = None
        
        while True:
            q = f"'{folder_id}' in parents and trashed = false"
            results = service.files().list(
                q=q, 
                pageSize=1000, 
                pageToken=page_token,
                orderBy=order_by,
                fields="nextPageToken, files(id, name, mimeType, modifiedTime, size, webViewLink, description)"
            ).execute()
            
            files = results.get('files', [])
            all_files.extend(files)
            
            page_token = results.get('nextPageToken')
            if not page_token:
                break
        
        return json.dumps(all_files)
    except Exception as e:
        return json.dumps({"error": str(e)})

@mcp.tool()
def search_files(token: str, query: str, order_by: str = 'folder,name') -> str:
    """
    Search for ALL matching files in Google Drive. Returns a JSON string.
    
    Args:
        token: The OAuth2 access token for the user.
        query: The search query (e.g., "name contains 'report'", "fullText contains 'budget'").
               Note: 'trashed = false' is automatically appended to the query.
        order_by: Sort order.
    """
    try:
        creds = Credentials(token=token)
        service = build('drive', 'v3', credentials=creds)

        all_files = []
        page_token = None
        
        # Enforce trashed = false if not present
        if "trashed" not in query:
            query = f"({query}) and trashed = false"

        while True:
            results = service.files().list(
                q=query, 
                pageSize=1000, 
                pageToken=page_token,
                orderBy=order_by,
                fields="nextPageToken, files(id, name, mimeType, modifiedTime, size, webViewLink, description)"
            ).execute()
            
            files = results.get('files', [])
            all_files.extend(files)
            
            page_token = results.get('nextPageToken')
            if not page_token:
                break
        
        return json.dumps(all_files)
    except Exception as e:
        return json.dumps({"error": str(e)})

@mcp.tool()
def read_file_content(token: str, file_id: str) -> str:
    """
    Read the content of a file from Google Drive.
    Supports text files, Google Docs, and PDFs.
    
    Args:
        token: The OAuth2 access token for the user.
        file_id: The ID of the file to read.
    """
    try:
        creds = Credentials(token=token)
        service = build('drive', 'v3', credentials=creds)

        # Get file metadata to check mimeType
        file_metadata = service.files().get(fileId=file_id, fields="id, name, mimeType").execute()
        mime_type = file_metadata.get('mimeType')
        file_name = file_metadata.get('name', '').lower()

        if mime_type == 'application/vnd.google-apps.document':
            # Export Google Doc to text
            request = service.files().export_media(fileId=file_id, mimeType='text/plain')
            file = io.BytesIO()
            downloader = MediaIoBaseDownload(file, request)
            done = False
            while done is False:
                status, done = downloader.next_chunk()
            return file.getvalue().decode('utf-8')
            
        elif mime_type == 'application/pdf' or file_name.endswith('.pdf'):
            # Download PDF and extract text
            request = service.files().get_media(fileId=file_id)
            file = io.BytesIO()
            downloader = MediaIoBaseDownload(file, request)
            done = False
            while done is False:
                status, done = downloader.next_chunk()
            
            file.seek(0)
            reader = PdfReader(file)
            text = ""
            for page in reader.pages:
                text += page.extract_text() + "\n"
            return text

        elif mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document' or file_name.endswith('.docx'):
            # Download DOCX and extract text
            request = service.files().get_media(fileId=file_id)
            file = io.BytesIO()
            downloader = MediaIoBaseDownload(file, request)
            done = False
            while done is False:
                status, done = downloader.next_chunk()
            
            file.seek(0)
            doc = Document(file)
            text = ""
            for para in doc.paragraphs:
                text += para.text + "\n"
            return text

        elif mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation' or file_name.endswith('.pptx'):
            # Download PPTX and extract text
            request = service.files().get_media(fileId=file_id)
            file = io.BytesIO()
            downloader = MediaIoBaseDownload(file, request)
            done = False
            while done is False:
                status, done = downloader.next_chunk()
            
            file.seek(0)
            prs = Presentation(file)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text

        elif mime_type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet' or file_name.endswith('.xlsx'):
            # Download XLSX and extract text
            request = service.files().get_media(fileId=file_id)
            file = io.BytesIO()
            downloader = MediaIoBaseDownload(file, request)
            done = False
            while done is False:
                status, done = downloader.next_chunk()
            
            file.seek(0)
            wb = openpyxl.load_workbook(file, data_only=True)
            text = ""
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                text += f"Sheet: {sheet}\n"
                for row in ws.iter_rows(values_only=True):
                    row_text = [str(cell) if cell is not None else "" for cell in row]
                    text += "\t".join(row_text) + "\n"
            return text

        else:
            # Assume text-based file (txt, json, csv, etc.)
            request = service.files().get_media(fileId=file_id)
            file = io.BytesIO()
            downloader = MediaIoBaseDownload(file, request)
            done = False
            while done is False:
                status, done = downloader.next_chunk()

            try:
                return file.getvalue().decode('utf-8')
            except UnicodeDecodeError:
                return f"Error: File type '{mime_type}' is not supported for text extraction, and could not be decoded as UTF-8 text."
    except Exception as e:
        return f"Error reading file: {str(e)}"

if __name__ == "__main__":
    mcp.run(transport="sse")
